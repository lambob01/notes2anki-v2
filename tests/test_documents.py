from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from notes2anki_v2.documents import ImageDocument, PptxDocument, open_document


def _blank_presentation() -> Presentation:
    presentation = Presentation()
    return presentation


def _add_blank_slide(presentation: Presentation):
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _make_png(tmp_path: Path) -> Path:
    image_path = tmp_path / "pic.png"
    Image.new("RGB", (10, 10), "red").save(image_path)
    return image_path


def test_open_document_dispatches_by_extension(tmp_path: Path) -> None:
    assert isinstance(open_document(tmp_path / "a.pptx"), PptxDocument)
    assert isinstance(open_document(tmp_path / "a.png"), ImageDocument)


def test_blank_slide_is_title_or_blank(tmp_path: Path) -> None:
    presentation = _blank_presentation()
    _add_blank_slide(presentation)
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    assert PptxDocument(pptx_path).is_title_or_blank(0) is True


def test_picture_slide_is_not_blank(tmp_path: Path) -> None:
    presentation = _blank_presentation()
    slide = _add_blank_slide(presentation)
    slide.shapes.add_picture(str(_make_png(tmp_path)), Inches(1), Inches(1))
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    assert PptxDocument(pptx_path).is_title_or_blank(0) is False


def test_text_heavy_slide_is_not_blank(tmp_path: Path) -> None:
    presentation = _blank_presentation()
    slide = _add_blank_slide(presentation)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    box.text_frame.text = "The Krebs cycle produces ATP through oxidative phosphorylation steps"
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    assert PptxDocument(pptx_path).is_title_or_blank(0) is False


def test_slide_notes_round_trip(tmp_path: Path) -> None:
    presentation = _blank_presentation()
    slide = _add_blank_slide(presentation)
    slide.notes_slide.notes_text_frame.text = "Remember the enzyme names"
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    document = PptxDocument(pptx_path)

    assert document.slide_notes(0) == "Remember the enzyme names"
    assert document.slide_notes(5) == ""


def test_presentation_is_parsed_once(tmp_path: Path) -> None:
    presentation = _blank_presentation()
    _add_blank_slide(presentation)
    pptx_path = tmp_path / "deck.pptx"
    presentation.save(pptx_path)

    document = PptxDocument(pptx_path)
    first = document.presentation
    document.slide_notes(0)
    document.is_title_or_blank(0)

    assert document.presentation is first
