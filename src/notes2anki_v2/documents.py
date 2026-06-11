from __future__ import annotations

import os
import shutil
import subprocess
from pathlib import Path
from typing import Callable

from PIL import Image, ImageDraw

MAX_CONTEXT_CHARS = 80_000
LIBREOFFICE_TIMEOUT_SECONDS = 180
MIN_CONTENT_WORDS = 6

WarnFn = Callable[[str], None]


class DocumentError(RuntimeError):
    pass


def open_document(path: Path) -> Document:
    extension = path.suffix.lower()
    if extension == ".pdf":
        return PdfDocument(path)
    if extension == ".pptx":
        return PptxDocument(path)
    return ImageDocument(path)


class Document:
    def __init__(self, path: Path) -> None:
        self.path = path

    def render_images(self, out_dir: Path, dpi: int, warn: WarnFn | None = None) -> list[Path]:
        raise NotImplementedError

    def slide_notes(self, index: int) -> str:
        return ""

    def is_title_or_blank(self, index: int) -> bool:
        return False

    def text_summary(self) -> str:
        return ""


class ImageDocument(Document):
    def render_images(self, out_dir: Path, dpi: int, warn: WarnFn | None = None) -> list[Path]:
        return [normalize_image_to_jpeg(self.path, out_dir)]


class PdfDocument(Document):
    def render_images(self, out_dir: Path, dpi: int, warn: WarnFn | None = None) -> list[Path]:
        return pdf_to_images(self.path, out_dir, dpi)

    def text_summary(self) -> str:
        try:
            import fitz  # type: ignore

            fitz.TOOLS.mupdf_display_errors(False)
            document = fitz.open(str(self.path))
            parts = []
            for page_number, page in enumerate(document, start=1):
                text = page.get_text().strip()
                if text:
                    parts.append(f"--- Slide {page_number} ---\n{text}")
            return "\n\n".join(parts)[:MAX_CONTEXT_CHARS]
        except Exception:
            return ""


class PptxDocument(Document):
    def __init__(self, path: Path) -> None:
        super().__init__(path)
        self._presentation = None

    @property
    def presentation(self):
        if self._presentation is None:
            try:
                from pptx import Presentation
            except ImportError as exc:
                raise DocumentError(
                    "python-pptx is not installed. Install dependencies with: pip install -e ."
                ) from exc
            try:
                self._presentation = Presentation(str(self.path))
            except Exception as exc:
                raise DocumentError(f"Could not open PPTX '{self.path.name}': {exc}") from exc
        return self._presentation

    def render_images(self, out_dir: Path, dpi: int, warn: WarnFn | None = None) -> list[Path]:
        fallback_reason = ""
        soffice = find_libreoffice()
        if soffice:
            try:
                subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        str(out_dir),
                        str(self.path),
                    ],
                    check=True,
                    capture_output=True,
                    timeout=LIBREOFFICE_TIMEOUT_SECONDS,
                )
                rendered_pdf = out_dir / f"{self.path.stem}.pdf"
                if rendered_pdf.exists():
                    images = pdf_to_images(rendered_pdf, out_dir, dpi)
                    if images:
                        return images
                fallback_reason = "LibreOffice did not produce a usable PDF"
            except subprocess.CalledProcessError as exc:
                stderr = (exc.stderr or b"").decode(errors="replace").strip()
                fallback_reason = f"LibreOffice conversion failed: {stderr or exc}"
            except Exception as exc:
                fallback_reason = f"LibreOffice conversion failed: {exc}"
        else:
            fallback_reason = "LibreOffice (soffice) was not found on this system"

        if warn:
            warn(
                f"{fallback_reason}. Rendering {self.path.name} as text-only images; "
                "images, charts, and layout will be missing. "
                "Install LibreOffice for full-quality slides."
            )
        return self._text_fallback_images(out_dir)

    def _text_fallback_images(self, out_dir: Path) -> list[Path]:
        try:
            images: list[Path] = []
            for index, slide in enumerate(self.presentation.slides):
                image = Image.new("RGB", (1280, 720), "white")
                draw = ImageDraw.Draw(image)
                y = 40
                for line in _slide_text_lines(slide):
                    draw.text((40, y), line[:180], fill="black")
                    y += 30
                    if y > 680:
                        break
                image_path = out_dir / f"slide-{index:03d}.jpg"
                image.save(image_path, "JPEG", quality=90)
                images.append(image_path)
            return images
        except DocumentError:
            raise
        except Exception as exc:
            raise DocumentError(f"Could not render PPTX '{self.path.name}': {exc}") from exc

    def slide_notes(self, index: int) -> str:
        try:
            slide = self.presentation.slides[index]
            if not slide.has_notes_slide:
                return ""
            frame = slide.notes_slide.notes_text_frame
            if frame is None:
                return ""
            return "\n".join(p.text.strip() for p in frame.paragraphs if p.text.strip())
        except Exception:
            return ""

    def is_title_or_blank(self, index: int) -> bool:
        try:
            slide = self.presentation.slides[index]
        except Exception:
            return False
        if _has_visual_content(slide):
            return False
        text_chunks = _slide_text_lines(slide)
        words = sum(len(chunk.split()) for chunk in text_chunks)
        if words < MIN_CONTENT_WORDS:
            return True

        has_non_title_shape = False
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                has_non_title_shape = True
                continue
            if getattr(shape.placeholder_format, "idx", None) not in (0, 1):
                has_non_title_shape = True
        return not has_non_title_shape and len(text_chunks) <= 2

    def text_summary(self) -> str:
        try:
            parts = []
            for index, slide in enumerate(self.presentation.slides, start=1):
                text = "\n".join(_slide_text_lines(slide))
                notes = self.slide_notes(index - 1)
                combined = text
                if notes:
                    combined = f"{combined}\n\n[Speaker Notes]: {notes}".strip()
                if combined:
                    parts.append(f"--- Slide {index} ---\n{combined}")
            return "\n\n".join(parts)[:MAX_CONTEXT_CHARS]
        except Exception:
            return ""


def find_libreoffice() -> str | None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice:
        return soffice
    if os.name == "nt":
        for candidate in (
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ):
            if Path(candidate).exists():
                return candidate
    return None


def pdf_to_images(pdf_path: Path, out_dir: Path, dpi: int) -> list[Path]:
    try:
        import fitz  # type: ignore
    except ImportError as exc:
        raise DocumentError(
            "PyMuPDF is not installed. Install dependencies with: pip install -e ."
        ) from exc

    fitz.TOOLS.mupdf_display_errors(False)
    images: list[Path] = []
    try:
        document = fitz.open(str(pdf_path))
        matrix = fitz.Matrix(dpi / 72.0, dpi / 72.0)
        for index, page in enumerate(document):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = out_dir / f"slide-{index:03d}.jpg"
            pixmap.save(str(image_path))
            images.append(image_path)
    except Exception as exc:
        raise DocumentError(f"Could not render PDF '{pdf_path.name}': {exc}") from exc
    return images


def normalize_image_to_jpeg(image_path: Path, out_dir: Path) -> Path:
    try:
        with Image.open(image_path) as image:
            converted = image.convert("RGB")
            out_path = out_dir / f"{image_path.stem}.jpg"
            converted.save(out_path, "JPEG", quality=95)
            return out_path
    except Exception as exc:
        raise DocumentError(f"Could not read image '{image_path.name}': {exc}") from exc


def _has_visual_content(slide: object) -> bool:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        return False

    visual_types = {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.GROUP,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
    }
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "shape_type", None) in visual_types:
            return True
    return False


def _slide_text_lines(slide: object) -> list[str]:
    lines: list[str] = []
    for shape in getattr(slide, "shapes", []):
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    lines.append(text)
    return lines
